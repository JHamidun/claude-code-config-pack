"""
Slide Image Generator — generates presentation slides as AI images via Gemini.
Uses Gemini Image (gemini-3.1-flash-image-preview / Nano Banana 2) for generation, packages into PPTX.

24 AI styles in 3 tiers:
- Manus Originals (7):  vinyl, whiteboard, grove, fresco, easel, diorama, chromatic
- Manus Hybrid (7):     sketch, glamour, amber, arctic, neon, patina, onyx
- Bonus (10):           chalkboard, notebook, blueprint, glassmorphism, corporate,
                        dark-tech, dashboard, infographic, watercolor, minimal-clean

13 HTML-only styles also available via slide_templates.py (cerulean, cobalt, etc.)

Usage:
    # From JSON config
    python whiteboard_generator.py generate slides.json ./output [style]

    # Quick single slide test
    python whiteboard_generator.py test "<prompt>" <output_dir> [style]

    # Regenerate a single slide from config
    python whiteboard_generator.py regenerate <config.json> <output_dir> <slide_id> [style]

    # Package existing images into PPTX
    python whiteboard_generator.py pptx <image_dir> <output.pptx>

    # Package images into PPTX with speaker notes
    python whiteboard_generator.py notes-pptx <image_dir> <output.pptx> <notes.json>

    # Create HTML preview
    python whiteboard_generator.py html <image_dir> <output.html>

    # List all available styles
    python whiteboard_generator.py styles
"""
import os
import sys
import time
import json
from pathlib import Path

# Remove conflicting env var before importing SDK
os.environ.pop('GEMINI_API_KEY', None)

from dotenv import load_dotenv
load_dotenv(os.path.expanduser("~/.claude/.credentials.master.env"))

from google import genai
from google.genai import types
from PIL import Image
import io

# Default image model = Nano Banana 2 (fast, cheap, great Cyrillic). Override for premium decks:
#   MANUS_SLIDES_MODEL=nano-banana-pro-preview   (Nano Banana Pro — richest detail + best incidental text)
#   MANUS_SLIDES_MODEL=gemini-3-pro-image        (Nano Banana Pro alias)
#   MANUS_SLIDES_MODEL=gemini-3.1-flash-lite-image (Nano Banana 2 Lite — cheaper)
# NOTE: gemini-3.5-flash is TEXT-ONLY (no image output) — do NOT use it here.
MODEL = os.getenv("MANUS_SLIDES_MODEL", "gemini-3.1-flash-image-preview")

client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))

# ============================================================
# STYLE PRESETS — Visual styles for AI-generated slides
# Each style is a prompt prefix that sets the visual context.
# The slide content prompt is appended after the prefix.
# ============================================================

STYLES = {
    # ================================================================
    # CUSTOM STYLES
    # ================================================================

    "exec-sketch": """A professional business presentation slide, 16:9, on a warm off-white cream background (#F1F3F5) that COMPLETELY covers all four corners — NO white rectangles, NO empty blocks, NO containers, NO reserved boxes in any corner. Clean editorial layout with generous whitespace. Crisp modern sans-serif typography (Manrope/Inter style) — titles bold dark indigo (#0B1021), body in dark gray, clear visual hierarchy. Content is structured in neat rounded white cards with thin borders and very soft shadows (corporate style). BUT all icons, diagrams, arrows, underlines and small illustrations are HAND-DRAWN in marker/sketch style with slightly imperfect lines, using indigo-blue (#3B5BDB) and copper-orange (#C77B30) markers — this blend of clean corporate typography with playful hand-drawn marker icons is the signature look. ALL TEXT IN RUSSIAN, perfectly readable Cyrillic (use Cyrillic letters И, С, Н, Р, К, Т — NOT Latin lookalikes I, C, H, P, K, T), no errors, NO duplicates, NO repeated words, NO truncated phrases. Render each text element EXACTLY ONCE as written.
The slide content:

""",

    "clean-marker": """A presentation slide on a pure clean white background. No room, no furniture, no whiteboard frame, no magnets, no chairs, no walls — just a clean white canvas.
Content is drawn in thick dry-erase marker style — hand-drawn, slightly imperfect lines, colorful marker strokes.
Use bold black marker for titles, blue marker for subtitles, and colored markers (red, green, orange, purple) for accents and highlights.
Icons and illustrations are sketched in marker style — simple, recognizable, slightly playful but professional.
Generous whitespace and breathing room between elements. Clean and uncluttered layout.
The hand-drawn marker content on this clean white slide:

""",

    # ================================================================
    # MANUS ORIGINALS — 7 prompt presets in the spirit of the image-based
    # Manus themes (the ones the Manus UI marks with a banana emoji)
    # ================================================================

    "vinyl": """A retro vintage poster-style presentation slide in the style of 1930s-1950s jazz/concert posters.
Bold geometric shapes and stylized silhouette illustrations of people and objects.
Limited but striking color palette: mustard yellow (#d4a017), deep burgundy/crimson (#8b1a1a), navy blue (#1a2744), cream/off-white (#f5f0e1).
Textured, slightly distressed paper background with subtle grain and aging effects.
Large bold display typography with dramatic size contrasts — headlines massive, subtitles elegant.
Art deco decorative elements: stars, swirls, musical notes, zigzag lines.
Each content slide maintains the same vintage poster aesthetic with illustrated characters and icons.
The content on this vintage poster slide:

""",

    "whiteboard": """A photorealistic image of a large whiteboard mounted on a wall in a modern office conference room.
The whiteboard has slight light reflections from overhead fluorescent lighting.
You can see blurred office chairs and a conference table at the bottom edge of the frame.
The whiteboard has a thin silver/gray aluminum frame.
At the bottom of the whiteboard there are 4-5 small colorful magnetic dots (blue, red, green, yellow).
Small marker tray at the bottom with colored dry-erase markers visible.
Previous erased content slightly visible as faint smudges in corners.
ON the whiteboard surface, drawn with thick dry-erase markers, the following content is displayed:

""",

    "grove": """An enchanted fairy-tale forest illustration style presentation slide.
Soft pastel watercolor palette: lavender (#c4b5d4), mint green (#b8d8c8), soft pink (#e8c4c8), pale yellow (#f0e6b8).
Whimsical storybook aesthetic with cute illustrated woodland animals (fox, owl, hedgehog, bunny).
Magical elements: tiny fairy lights glowing on tree branches, sparkling dust paths, luminescent mushrooms.
Twisted ancient trees forming natural frames around the content area.
Hand-drawn organic lettering with swirls and flourishes, warm brown ink color.
Dreamy, ethereal atmosphere with soft gradients and translucent layers.
Visible paper/parchment texture underneath the watercolor washes.
The enchanting content on this slide:

""",

    "fresco": """A warm urban illustration style presentation slide, like a city guidebook or travel poster.
Flat design city scenes with buildings, people walking, cyclists, trees, cafes, bridges, trams.
Muted warm palette: terracotta (#c2704f), sage green (#8faa7b), warm beige (#e8dcc8), soft brown (#8b6f4f), dusty rose (#c4908a).
Hand-painted brush stroke texture, slightly loose and artistic.
Whimsical characters of diverse people going about daily city life.
Slightly vintage European city feel — cobblestone streets, arched windows, awnings.
Content slides use the same illustrated city elements and warm organic tones.
Title text in elegant warm gold/brown handwritten style.
The urban story on this slide:

""",

    "easel": """An artist's studio/atelier style presentation slide.
Depicted as a creative workspace scene: canvas on wooden easel, paint palettes with mixed colors, scattered brushes and pencils, open sketchbooks with botanical drawings.
Rich warm color palette: terracotta (#b85c38), warm brown (#6b4423), amber (#c68642), cream (#f5eed5), deep green (#3a5a40).
Warm golden lighting from a nearby window, creating cozy atmosphere.
Books stacked on shelves, potted plants, a steaming cup of coffee on the desk.
Cork board in background with pinned artwork, color swatches, and reference photos.
Content is presented as if on canvases, sketchbook pages, or pinned notes on the cork board.
Rich illustrated style with visible brushwork and artistic textures.
The creative content on this slide:

""",

    "diorama": """A pop-up book / paper craft style presentation slide.
3D paper elements rising dramatically from the pages of an open red hardcover book lying on a dark teal surface.
Paper-cut aesthetic with visible layered depth, shadows between paper layers.
Bright cheerful color palette: teal (#0d9488), lime green (#65a30d), coral (#f97316), sunny yellow (#facc15), sky blue (#38bdf8).
Buildings, trees, windmills, vehicles, people rendered as colorful paper cutouts with folded edges.
Slightly isometric 3D perspective showing the pop-up depth.
Decorative banner/ribbon elements for title text, floating above the scene.
Content slides maintain the pop-up book effect with paper-craft infographics and charts.
The pop-up content on this slide:

""",

    "chromatic": """A colorful, playful tech-explainer style presentation slide.
Light warm background: soft cream/peach (#fdf2e9) with subtle flowing curved lines in pastel.
Floating 3D rendered objects relevant to the topic — glossy, colorful, with reflections and soft shadows.
Rainbow prismatic color effects: light refracting through prisms, color spectrum arcs.
Playful but informative visual style — educational and engaging.
Large coral/salmon (#e8734a) headline text, bold and rounded sans-serif.
Blue (#4a7ce8) accent text for subtitles and secondary info.
Sparkles, translucent bubbles, flowing gradient ribbons as decorative elements.
Geometric shapes (spheres, cubes, diamonds) floating around content.
Content slides use 3D illustrated objects and icons to visually explain concepts.
The colorful content on this slide:

""",

    # ================================================================
    # MANUS "HTML" STYLES — visual look reproduced here via Gemini image generation
    # (in this pipeline every style below is rendered as an image, not as CSS)
    # ================================================================

    "sketch": """A hand-drawn chalk and charcoal sketch style presentation slide on a dark charcoal/black paper background.
Everything is drawn in white chalk, charcoal, and pencil with a hand-sketched imperfect quality.
Whimsical stick figures and cartoon characters with expressive poses and movements.
Bold white hand-lettered title text, slightly wobbly and organic.
Decorative swirls, stars, smoke wisps, and doodle elements scattered around.
Scientific diagrams drawn in sketch style with labels and arrows.
The artistic style resembles a talented artist's chalkboard doodles — playful but detailed.
Visible chalk dust texture and smudging effects on the dark background.
The sketched content on this dark board:

""",

    "glamour": """A luxury high-fashion editorial presentation slide.
Rich dark brown/black cinematic background with dramatic moody lighting.
Photorealistic flowing fabric (silk, satin) in warm golden/champagne tones, billowing gracefully.
Massive bold white serif typography (like Playfair Display or Didot) for the headline.
Elegant thin letter-spaced subtitle text in light gray.
Fashion photography aesthetic — dramatic shadows, warm golden highlights.
Premium, sophisticated feel. Minimal text, maximum visual impact.
Small footer text with year/date in the bottom left.
The luxurious content on this slide:

""",

    "amber": """A soft, organic abstract presentation slide with a calm, mindful aesthetic.
Light warm background with abstract flowing organic shapes in muted pastel tones:
soft beige (#e8dcc8), dusty blue (#b8c8d8), warm gold/amber (#d4a86a), sage (#c8d4c0).
Shapes are rounded, blob-like, overlapping with subtle transparency.
Clean modern sans-serif typography in dark charcoal gray (#333333).
Large centered headline text. Generous whitespace and breathing room.
Subtle dot patterns and organic textures within the abstract shapes.
Calming, professional, wellness-oriented aesthetic.
The serene content on this slide:

""",

    "arctic": """A cool corporate tech presentation slide with a photorealistic technology background.
Light silver-gray background with a subtle, blurred photograph of technology (robot hands, circuits, screens).
The photo is desaturated and low-contrast, serving as a muted backdrop.
Clean dark text overlaid on the light background with good contrast.
Blue (#2563eb) accent color for highlights and links.
Professional, understated, modern tech company aesthetic.
Small date or category text in the top left.
Large bold sans-serif headline. Structured content below.
The tech-forward content on this slide:

""",

    "neon": """A vibrant neon-themed presentation slide with a dark background.
Deep dark purple/black (#1a0a2e) background with electric neon elements.
Glowing neon text effects — hot pink (#ff1493), electric blue (#00d4ff), neon green (#39ff14), bright yellow (#ffff00).
Retro 80s/synthwave aesthetic with neon sign tube lettering for headlines.
Subtle grid perspective floor fading into the distance.
Neon light glow and bloom effects around key elements.
Bold, energetic, nightclub/gaming aesthetic.
The neon-lit content on this slide:

""",

    "patina": """A prehistoric cave painting / ancient art style presentation slide.
Warm ochre/terracotta cave wall texture as background (#c2874a, #8b6542).
Figures and symbols drawn in the style of ancient cave paintings — handprints, animals, stick figures.
Earthy color palette: red ochre (#a0522d), charcoal black, cream white, dusty orange.
Title text styled as if carved or painted on stone — rough, primitive but legible.
Archaeological/anthropological aesthetic. Textured stone surface visible.
The ancient content on this slide:

""",

    "onyx": """A bold philosophical/existential dark presentation slide.
Pure black (#000000) background with stark white text.
Massive bold condensed sans-serif headline text filling most of the frame.
Ultra-minimalist — almost brutalist typography design.
No decorations, no images — just powerful text on black.
Subtle dark gray (#1a1a1a) geometric element or circle behind the text.
Art gallery / museum exhibition poster aesthetic.
The bold statement on this slide:

""",

    # ================================================================
    # MANUS 1.6 IMAGE THEMES (added 2026-07-02)
    # Prompt presets written in the spirit of the Manus 1.6 image themes,
    # composed by looking at published sample slides for each theme.
    # Reference sample slides: references/manus-theme-samples/<Theme>__gpt-image/
    # ================================================================

    "etching": """A refined pen-and-ink engraving / etching style presentation slide on a warm off-white cream paper background (#F4F1EA) that FULLY covers all four corners — NO white rectangles, NO empty boxes, NO reserved containers in any corner.
Delicate fine-line ink illustration with subtle crosshatching and a soft pale grey watercolor wash, in the elegant restrained manner of a New Yorker illustration or a vintage encyclopedia engraving.
Title in a classic letterpress serif with generous letter-spacing (spaced roman capitals, Caslon/Garamond feel), thin dark charcoal ink.
Small understated italic serif subtitle beneath, often framed between two short em-dashes.
Muted monochrome palette: warm cream paper, charcoal/graphite ink, pale grey washes — no bright colors.
Extremely refined, quiet, literary and editorial, with very generous whitespace and a small italic caption or roman-numeral page mark near the bottom.
ALL TEXT IN RUSSIAN, perfectly readable Cyrillic (use Cyrillic letters И, С, Н, Р, К, Т — NOT Latin lookalikes I, C, H, P, K, T), no errors, NO duplicates, render each text element EXACTLY ONCE.
The refined engraved content on this slide:

""",

    "editorial": """A premium print-magazine editorial presentation slide in the sophisticated style of Kinfolk or Monocle, with TWO layout variants — pick by the slide content below:

IF the content below is a TITLE or SECTION slide: two-column split layout.
LEFT column: a moody, warm, photorealistic editorial photograph lit by soft natural window light — tactile objects (stacked design books, a ceramic coffee cup, dried flowers, textured linen, a film camera, loose paper sheets — vary the objects), gentle shadows.
RIGHT column: a solid dark charcoal/espresso panel (#211E1B). On it: a large high-contrast Didone serif display headline (Playfair Display style) in warm cream/ivory stacked over several lines; a single thin copper-orange rule under the headline; a refined serif subtitle in soft warm grey; at the bottom a small letter-spaced ALL-CAPS footer line in muted copper.

OTHERWISE (lists, diagrams, comparisons, data — any content slide): a light editorial "paper" layout on warm cream (#F4F1EA) covering the WHOLE slide, no dark panel.
Headline in the same dark charcoal Didone serif, moderate size (2–3 lines max, upper-left, avoid hyphenating or breaking words), with a thin copper-orange rule beneath. ONLY if the content below explicitly provides a subtitle or kicker, set it in italic serif under the rule; if NO subtitle is given, leave that area empty — NEVER invent a kicker or subtitle, never render instruction words.
Body set in an elegant editorial grid with generous whitespace: numbered lists with bold serif lead-ins separated by thin hairlines; comparisons and diagrams as columns divided by a thin vertical hairline, with large Didone numerals or copper accents for key figures and arrows in fine charcoal line; optionally a small editorial photograph inside a thin-bordered frame on one side.
Footer: a small letter-spaced ALL-CAPS line in muted copper at the bottom, above or below a thin hairline.

BOTH variants: a thin charcoal border frame around the whole slide; palette strictly warm cream, dark charcoal/espresso, muted copper-orange; understated luxury print aesthetic; body text clearly smaller than the headline, headline never taking more than a third of the layout on content slides.
ALL TEXT IN RUSSIAN, perfectly readable Cyrillic (use Cyrillic letters И, С, Н, Р, К, Т — NOT Latin lookalikes I, C, H, P, K, T), no errors, NO duplicates, render each text element EXACTLY ONCE.
The sophisticated editorial content on this slide:

""",

    "pixel": """A nostalgic late-1990s / early-2000s classic Mac OS desktop-GUI style presentation slide.
The entire slide is one beveled 'Platinum'/Aqua application window: a light grey pinstriped title bar with a short centered label (use the window label given in the content; if none is given, label it "Presentation") and three traffic-light dots (red, yellow, green) in the top-left corner.
The area around the window is a bright Bondi-blue retro desktop with a subtle dotted / pinstripe halftone texture.
Inside the window, on an off-white panel: a big, bold, slightly rounded grotesque headline (Charcoal/Chicago-era feel) in near-black, tightly set across two or three lines, with a short orange underline accent beneath it.
If the content provides a subtitle, set it in a lighter grey sans-serif below the headline; if NO subtitle is given, leave that area empty — NEVER invent one.
Footer: ONLY if the content explicitly lists footer tags, render a recessed inset status/footer bar at the bottom holding exactly those bullet-separated tags in bold or monospaced text. If NO footer tags are given, OMIT the footer bar entirely — never invent tag words, never reuse words from the body as tags.
Skeuomorphic beveled edges, soft inset shadows, warm retro-computing nostalgia; optionally a faint blue line-art of a vintage gadget in one corner.
Render ONLY the text given in the content below — no extra labels, captions or filler words of your own.
ALL TEXT IN RUSSIAN, perfectly readable Cyrillic (use Cyrillic letters И, С, Н, Р, К, Т — NOT Latin lookalikes I, C, H, P, K, T), no errors, NO duplicates, render each text element EXACTLY ONCE.
The retro-computing content on this slide:

""",

    "vellum": """A serene East-Asian ink-wash (sumi-e / shan-shui) style presentation slide on a warm rice-paper cream background (#F2ECDD) with a subtle paper texture that covers all corners.
A soft monochrome ink-wash painting of misty, layered mountains with a gnarled pine tree occupies the lower portion, in muted grey-green and charcoal ink with gentle bleeds and plenty of negative space; on dense content slides (lists, comparisons, diagrams) keep the landscape within roughly the bottom 25%, faint and semi-transparent, so it NEVER touches or sits behind the text.
Title in an elegant serif (refined Song/Ming-inspired or classic Western serif) in deep ink-green (#2E4034), calm and balanced.
If the content provides a subtitle, set it in a restrained fine serif or sans-serif beneath the title; if NO subtitle is given, leave that area empty — NEVER invent one.
ALWAYS include exactly one small red seal / chop stamp (vermilion #B23A2E) — the ONLY bright accent color, present on every slide: place it quietly right beneath the title block or in the bottom-right corner — never floating in the middle of the content area.
Zen, meditative, classical, spacious; all text sits in the clear upper area above the mountains.
Render ONLY the text given in the content below — no extra labels, captions or filler words of your own.
ALL TEXT IN RUSSIAN, perfectly readable Cyrillic (use Cyrillic letters И, С, Н, Р, К, Т — NOT Latin lookalikes I, C, H, P, K, T), no errors, NO duplicates, render each text element EXACTLY ONCE.
The tranquil content on this slide:

""",

    "dossier": """A vintage archival 'dossier' flat-lay presentation slide — a top-down photorealistic still-life.
The centerpiece is a large aged, slightly torn cream paper note (faint ruled lines, soft drop shadow, held by a small brass pin) resting on a warm textured cork or tan-leather desk surface that fills the whole frame.
Curated antique objects are arranged ONLY around the outer edges of the frame — a brass service bell, old postage stamps, a wax-seal stamp with a small red wax seal, an antique brass key on a fob, white gloves, a fountain pen, a wax-sealed envelope — all in warm brass, sepia and cream tones; objects may be partly cropped by the frame but must NEVER overlap or touch the paper's text. ALL text sits in the clean center of the aged paper.
On the aged paper: the title in a classic serif (Playfair / Georgia feel) in deep navy (#1F3A5F), with small ornamental serif dividers above or below it.
If the content provides a subtitle, set it in a clean serif beneath the title; if NO subtitle is given, leave that area empty — NEVER invent one.
If the content provides a short red accent line (e.g. an "Est. ..." year), render it in small red italic serif between two fine flourishes; if NO red accent line is given, OMIT it entirely — never invent dates, mottos or slogans.
On dense content slides (lists, comparisons, diagrams): list items get small brass ornament markers and are separated by thin sepia hairlines; comparisons form two columns divided by a thin vertical sepia hairline; body text in dark ink serif, clearly smaller than the title.
Warm brown/sepia lighting with rich soft shadows and tactile realism — an elegant heritage / detective-dossier aesthetic.
Render ONLY the text given in the content below — no extra labels, captions, worded stamps or filler words of your own.
ALL TEXT IN RUSSIAN, perfectly readable Cyrillic (use Cyrillic letters И, С, Н, Р, К, Т — NOT Latin lookalikes I, C, H, P, K, T), no errors, NO duplicates, render each text element EXACTLY ONCE.
The archival content on this slide:

""",

    "sketch-notebook": """A friendly hand-drawn ink doodle presentation slide on a warm off-white cream notebook paper background (#F7F4EC) with a very faint light-grey graph-paper grid, covering the whole frame including all corners.
Everything is drawn in thin black ink pen / fine felt-tip with a charming, slightly imperfect hand-drawn quality — pure monochrome black ink on cream, NO other colors.
Big friendly rounded hand-lettered headline in black ink: soft rounded letterforms like careful casual handwriting, generously sized.
If the content provides a subtitle, hand-letter it smaller beneath or below the main drawing; if NO subtitle is given, leave that area empty — NEVER invent one.
Cute simple ink-line cartoon people (thin outlines, simple dot-and-line faces, light scribble hatching on hair and trousers) interacting with the content — pointing, presenting, walking.
Simple doodle icons and props drawn in the same thin ink line (coffee cup, clock, arrows, picture frames, classical columns as pedestals for list items), wavy hand-drawn underlines and dividers.
A few loose decorative spiral doodles in the corners. Generous whitespace, light and airy, warm and approachable — like a beautifully doodled Moleskine page or a friendly explainer comic.
For lists and diagrams: each item gets its own little doodle icon, hand-written labels beneath, connected by sketchy arrows.
Render ONLY the text given in the content below — no extra labels, captions or filler words of your own.
ALL TEXT IN RUSSIAN, perfectly readable Cyrillic (use Cyrillic letters И, С, Н, Р, К, Т — NOT Latin lookalikes I, C, H, P, K, T), no errors, NO duplicates, render each text element EXACTLY ONCE.
The hand-doodled content on this slide:

""",

    # ================================================================
    # BONUS STYLES — Additional original styles
    # ================================================================

    "chalkboard": """A photorealistic image of a large green chalkboard (blackboard) in a university lecture hall.
Visible chalk dust and smudges. Thick wooden frame. Chalk tray at the bottom with white and colored chalk.
The chalk writing has a natural, slightly imperfect hand-drawn quality.
Warm incandescent lighting from above. The lecture hall is dimly visible in the background.
Written in chalk on the board:

""",

    "notebook": """A photorealistic top-down photograph of an open Moleskine-style notebook on a wooden desk.
Cream/off-white lined paper. Content is written with a fine-tip blue and black pen.
Hand-drawn diagrams with ruler-straight lines. Highlighted sections with yellow and pink markers.
Small colorful sticky notes attached to the page edges. A pen lies across the bottom of the page.
Written and drawn in the notebook:

""",

    "blueprint": """A photorealistic rendering of a technical blueprint-style presentation slide.
Deep blue (#003366) background with white grid lines forming a precise engineering grid.
All content rendered in white and light cyan lines, mimicking technical drawings.
Precise geometric shapes, measurement annotations, technical labels.
Title block in the bottom right corner with date and revision number.
Monospace typography throughout. Engineering/architecture aesthetic.
The technical content on this blueprint:

""",

    "glassmorphism": """A modern glassmorphism-style presentation slide.
Vibrant gradient background with purple (#8b5cf6), blue (#3b82f6), and pink (#ec4899) blobs.
Content is displayed on frosted glass panels with blur effect and subtle white borders.
Semi-transparent cards floating over the colorful background.
Clean white text on the glass panels. Smooth rounded corners. Modern, premium aesthetic.
The content on the glass panels:

""",

    "corporate": """A professional corporate presentation slide displayed on a large screen.
Clean white background with a navy blue (#1a365d) header bar at the top.
Professional sans-serif typography. Structured grid layout.
Subtle shadow effects on content cards. Blue accent color for highlights.
Page number in the bottom right corner.
The content on this slide is:

""",

    "dark-tech": """A dark tech-style presentation slide.
Deep black (#0a0a0a) background with subtle dark grid pattern.
Neon green (#00ff88) and cyan (#00d4ff) accent lines and highlights.
Monospace font for technical content, clean sans-serif for headers.
Subtle glow effects on key elements. Terminal/hacker aesthetic with professional polish.
The content on this slide is:

""",

    "dashboard": """A data dashboard presentation slide on a dark screen.
Deep navy (#0f172a) to dark gray (#1e293b) background.
Content organized in card-based layout with subtle borders and rounded corners.
Charts in blue (#3b82f6), green (#22c55e), orange (#f59e0b), red (#ef4444).
Large bold numbers for metrics. Clean sans-serif typography.
The data on this slide is:

""",

    "infographic": """A colorful infographic-style slide with a clean white background.
Flat design icons and illustrations. Bold, large numbers for statistics.
Vibrant palette: blue (#2563eb), teal (#0d9488), orange (#ea580c), purple (#9333ea).
Content in clear sections with visual hierarchy. Rounded shapes and smooth lines.
Data visualizations: bar charts, pie charts, progress bars, icon arrays.
The content on this slide is:

""",

    "watercolor": """An artistic presentation slide rendered in watercolor painting style.
Soft, flowing watercolor washes as background — pale blues, soft pinks, muted greens.
Text in elegant dark gray calligraphic style over the watercolor.
Illustrations painted in loose watercolor style with visible brush strokes.
Delicate, artistic aesthetic. White paper texture visible underneath.
The content painted on this slide:

""",

    "minimal-clean": """An ultra-minimalist presentation slide with maximum whitespace.
Pure white (#ffffff) background. Content centered with extreme precision.
Black (#000000) text only. Thin, elegant sans-serif font.
No decorations, no icons, no borders — pure typography and whitespace.
A single thin hairline divider if needed. Less is more.
The content on this slide:

""",
}

# Style aliases for convenience
STYLE_ALIASES = {
    # exec-sketch aliases (business deck + hand-drawn marker icons)
    "yourname": "exec-sketch",
    "business-sketch": "exec-sketch",
    "exec": "exec-sketch",
    # Whiteboard aliases
    "office": "whiteboard",
    "marker": "whiteboard",
    "marker-board": "whiteboard",
    # Manus AI style aliases
    "retro": "vinyl",
    "poster": "vinyl",
    "jazz": "vinyl",
    "fairy": "grove",
    "fantasy": "grove",
    "storybook": "grove",
    "urban": "fresco",
    "city": "fresco",
    "travel": "fresco",
    "art-studio": "easel",
    "atelier": "easel",
    "craft": "easel",
    "popup": "diorama",
    "paper": "diorama",
    "book": "diorama",
    "rainbow": "chromatic",
    "playful": "chromatic",
    "explainer": "chromatic",
    # Manus "HTML" style aliases
    "doodle": "sketch",
    "chalk-sketch": "sketch",
    "fashion": "glamour",
    "luxury": "glamour",
    "haute": "glamour",
    "calm": "amber",
    "wellness": "amber",
    "organic": "amber",
    "tech-photo": "arctic",
    "robot": "arctic",
    "synthwave": "neon",
    "80s": "neon",
    "cave": "patina",
    "ancient": "patina",
    "prehistoric": "patina",
    "brutalist": "onyx",
    "existential": "onyx",
    "black": "onyx",
    # Bonus style aliases
    "chalk": "chalkboard",
    "school": "chalkboard",
    "handwritten": "notebook",
    "journal": "notebook",
    "technical": "blueprint",
    "engineering": "blueprint",
    "modern": "glassmorphism",
    "glass": "glassmorphism",
    "professional": "corporate",
    "business": "corporate",
    "tech": "dark-tech",
    "hacker": "dark-tech",
    "data": "dashboard",
    "metrics": "dashboard",
    "stats": "infographic",
    "artistic": "watercolor",
    "paint": "watercolor",
    "clean": "minimal-clean",
    "minimal": "minimal-clean",
    # Manus 1.6 image-theme aliases (added 2026-07-02)
    "engraving": "etching",
    "ink-engraving": "etching",
    "newyorker": "etching",
    "magazine": "editorial",
    "kinfolk": "editorial",
    "monocle": "editorial",
    "retro-mac": "pixel",
    "macos": "pixel",
    "aqua": "pixel",
    "y2k": "pixel",
    "ink-wash": "vellum",
    "sumi": "vellum",
    "sumi-e": "vellum",
    "zen": "vellum",
    "shanshui": "vellum",
    "archival": "dossier",
    "flatlay": "dossier",
    "flat-lay": "dossier",
    "vintage-desk": "dossier",
    "heritage": "dossier",
    "manus-sketch": "sketch-notebook",
    "ink-doodle": "sketch-notebook",
    "notebook-sketch": "sketch-notebook",
}

# Category groupings for display
STYLE_CATEGORIES = {
    "Manus AI (Nano Banana originals)": [
        "vinyl", "whiteboard", "grove", "fresco", "easel", "diorama", "chromatic",
    ],
    "Manus Hybrid (rendered via Gemini)": [
        "sketch", "glamour", "amber", "arctic", "neon", "patina", "onyx",
    ],
    "Bonus AI styles": [
        "chalkboard", "notebook", "blueprint", "glassmorphism",
        "corporate", "dark-tech", "dashboard", "infographic",
        "watercolor", "minimal-clean",
    ],
    "Manus 1.6 image themes (prompt presets)": [
        "etching", "editorial", "pixel", "vellum", "dossier", "sketch-notebook",
    ],
}

# Manus HTML-only templates (pure CSS/typography, use slide_templates.py for these)
# These DON'T need AI generation — they are clean CSS layouts
HTML_ONLY_STYLES = [
    "cerulean", "cobalt", "emerald", "basalt",
    "mist", "sand", "linen", "alabaster",
    "quartz", "mahogany", "ginkgo", "sunset", "lavender",
]

# Recommendation engine: task type -> best style
STYLE_RECOMMENDATIONS = {
    "pitch deck": ["vinyl", "diorama", "chromatic"],
    "investor deck": ["corporate", "vinyl", "dashboard"],
    "startup": ["chromatic", "glassmorphism", "diorama"],
    "education": ["whiteboard", "chalkboard", "sketch", "grove"],
    "children": ["grove", "diorama", "chromatic", "sketch"],
    "tech": ["dark-tech", "arctic", "blueprint", "dashboard"],
    "data": ["dashboard", "infographic", "corporate"],
    "creative": ["easel", "watercolor", "fresco", "sketch"],
    "art": ["easel", "watercolor", "vinyl", "sketch"],
    "travel": ["fresco", "vinyl", "grove"],
    "science": ["sketch", "blueprint", "chalkboard", "infographic"],
    "conference": ["whiteboard", "corporate", "vinyl"],
    "internal": ["whiteboard", "corporate", "minimal-clean"],
    "luxury": ["glamour", "minimal-clean", "glassmorphism", "onyx"],
    "fashion": ["glamour", "vinyl", "onyx"],
    "fun": ["chromatic", "diorama", "grove", "sketch"],
    "wellness": ["amber", "grove", "watercolor"],
    "philosophy": ["onyx", "amber", "minimal-clean"],
    "history": ["patina", "vinyl", "fresco"],
    "gaming": ["neon", "dark-tech", "chromatic"],
    "music": ["vinyl", "neon", "sketch"],
    "event": ["neon", "chromatic", "diorama"],
    # Manus 1.6 image themes (added 2026-07-02)
    "magazine": ["editorial", "glamour", "minimal-clean"],
    "editorial": ["editorial", "etching", "glamour"],
    "brand": ["editorial", "glamour", "vinyl"],
    "hospitality": ["dossier", "editorial", "patina"],
    "heritage": ["dossier", "patina", "vellum"],
    "mindfulness": ["vellum", "amber", "watercolor"],
    "zen": ["vellum", "amber", "minimal-clean"],
    "literary": ["etching", "editorial", "minimal-clean"],
    "retro": ["pixel", "vinyl", "neon"],
    "nostalgia": ["pixel", "vinyl", "patina"],
    "product design": ["pixel", "editorial", "arctic"],
}

DEFAULT_STYLE = "whiteboard"


def resolve_style(style: str) -> str:
    """Resolve style name through aliases. Returns canonical style name."""
    resolved = STYLE_ALIASES.get(style, style)
    if resolved in STYLES:
        return resolved
    return DEFAULT_STYLE


def recommend_styles(task_description: str) -> list:
    """Recommend best styles for a given task description."""
    task_lower = task_description.lower()
    matches = []
    for keyword, styles in STYLE_RECOMMENDATIONS.items():
        if keyword in task_lower:
            matches.extend(styles)
    # Deduplicate preserving order
    seen = set()
    result = []
    for s in matches:
        if s not in seen:
            seen.add(s)
            result.append(s)
    return result[:5] if result else ["whiteboard", "vinyl", "chromatic"]


def generate_slide_image(
    prompt: str,
    output_path: Path,
    style: str = DEFAULT_STYLE,
    width: int = 1280,
    height: int = 720,
    retry_count: int = 3,
    temperature: float = 0.4,
) -> bool:
    """Generate a single slide image via Gemini.

    Args:
        prompt: Content description (what to draw/write on the slide)
        output_path: Where to save the image
        style: Visual style preset name (see STYLES dict)
        width: Target width in pixels
        height: Target height in pixels
        retry_count: Number of retries on failure
        temperature: Generation temperature (0.0-1.0)

    Returns:
        True if successful, False otherwise
    """
    resolved = resolve_style(style)
    style_prefix = STYLES[resolved]
    full_prompt = style_prefix + prompt

    for attempt in range(retry_count):
        try:
            response = client.models.generate_content(
                model=MODEL,
                contents=full_prompt,
                config=types.GenerateContentConfig(
                    response_modalities=['IMAGE', 'TEXT'],
                    temperature=temperature,
                )
            )

            if not response.candidates or not response.candidates[0].content:
                print(f"  [WARN] Empty response (attempt {attempt + 1})")
                if attempt < retry_count - 1:
                    time.sleep(3)
                continue

            for part in response.candidates[0].content.parts:
                if part.inline_data and part.inline_data.mime_type and part.inline_data.mime_type.startswith("image/"):
                    img = Image.open(io.BytesIO(part.inline_data.data))
                    if img.size != (width, height):
                        img = img.resize((width, height), Image.LANCZOS)
                    output_path.parent.mkdir(parents=True, exist_ok=True)
                    img.save(str(output_path), "PNG")
                    return True

            print(f"  [WARN] No image in response (attempt {attempt + 1})")
            if attempt < retry_count - 1:
                time.sleep(3)

        except Exception as e:
            print(f"  [ERR] Attempt {attempt + 1}: {e}")
            if attempt < retry_count - 1:
                time.sleep(5)

    return False


def generate_from_config(config_path: str, output_dir: str, style: str = DEFAULT_STYLE) -> list:
    """Generate all slides from a JSON config file.

    Config format:
    {
        "title": "Presentation Title",
        "style": "vinyl",  // optional, overrides default
        "slides": [
            {"id": "slide_01", "prompt": "Content description..."},
            ...
        ]
    }

    Returns list of generated image paths.
    """
    config = json.loads(Path(config_path).read_text(encoding="utf-8"))
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)

    # CLI style overrides config style (if CLI provides non-default)
    if style != DEFAULT_STYLE:
        slide_style = style
    else:
        slide_style = config.get("style", style)
    slides = config["slides"]
    image_paths = []

    print(f"Generating {len(slides)} slides (style: {slide_style})...")

    for i, slide in enumerate(slides, 1):
        slide_id = slide.get("id", f"slide_{i:02d}")
        img_path = out / f"{slide_id}.png"

        if img_path.exists():
            print(f"  [{i}/{len(slides)}] {slide_id} — already exists, skipping")
            image_paths.append(img_path)
            continue

        print(f"  [{i}/{len(slides)}] {slide_id} — generating...")
        ok = generate_slide_image(slide["prompt"], img_path, style=slide_style)

        if ok:
            print(f"  [{i}/{len(slides)}] {slide_id} — OK")
            image_paths.append(img_path)
        else:
            print(f"  [{i}/{len(slides)}] {slide_id} — FAILED")
            image_paths.append(None)

        # Rate limiting
        if i < len(slides):
            time.sleep(2)

    success = sum(1 for p in image_paths if p)
    print(f"\nGenerated: {success}/{len(slides)} slides")
    return image_paths


def create_pptx(image_dir: str, output_path: str, slide_width_inches: float = 13.333, slide_height_inches: float = 7.5):
    """Package slide images into PPTX."""
    from pptx import Presentation
    from pptx.util import Inches, Emu

    img_dir = Path(image_dir)
    images = sorted(img_dir.glob("slide_*.png"))

    if not images:
        print("No slide images found!")
        return False

    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)

    blank_layout = prs.slide_layouts[6]

    for img_path in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )

    prs.save(str(output_path))
    print(f"PPTX saved: {output_path} ({len(images)} slides)")
    return True


def create_html_preview(image_dir: str, output_path: str):
    """Create navigable HTML preview of slides."""
    img_dir = Path(image_dir)
    out_path = Path(output_path)
    images = sorted(img_dir.glob("slide_*.png"))

    if not images:
        print("No slide images found!")
        return False

    slides_html = ""
    for i, img in enumerate(images):
        rel = img.relative_to(out_path.parent)
        slides_html += f"""
        <div class="slide{' active' if i == 0 else ''}" data-idx="{i}">
            <img src="{rel}" alt="Slide {i+1}">
            <div class="num">{i+1}/{len(images)}</div>
        </div>"""

    html = f"""<!DOCTYPE html>
<html lang="ru">
<head>
<meta charset="UTF-8">
<title>Slide Presentation</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#1a1a1a;font-family:Arial}}
.slide{{display:none;width:100vw;height:100vh;justify-content:center;align-items:center;position:relative}}
.slide.active{{display:flex}}
.slide img{{max-width:95vw;max-height:95vh;object-fit:contain;border-radius:4px;box-shadow:0 4px 20px rgba(0,0,0,.5)}}
.num{{position:absolute;bottom:10px;right:20px;color:#888;font-size:14px}}
.hint{{position:fixed;bottom:10px;left:50%;transform:translateX(-50%);color:#555;font-size:12px}}
</style>
</head>
<body>
{slides_html}
<div class="hint">Arrow keys or click to navigate</div>
<script>
let c=0;const s=document.querySelectorAll('.slide');
function go(n){{s[c].classList.remove('active');c=(n+s.length)%s.length;s[c].classList.add('active')}}
document.addEventListener('keydown',e=>{{if(e.key==='ArrowRight'||e.key===' ')go(c+1);if(e.key==='ArrowLeft')go(c-1)}});
document.addEventListener('click',e=>{{e.clientX>innerWidth/2?go(c+1):go(c-1)}});
</script>
</body>
</html>"""

    out_path.write_text(html, encoding="utf-8")
    print(f"HTML preview saved: {output_path}")
    return True


def regenerate_slide(config_path: str, output_dir: str, slide_id: str, style: str = DEFAULT_STYLE) -> bool:
    """Regenerate a single slide from config, then rebuild PPTX + HTML.

    Args:
        config_path: Path to slides.json config file.
        output_dir: Directory containing slide images.
        slide_id: Slide ID to regenerate (e.g. "slide_01").
        style: Visual style override (uses config style if not specified).

    Returns:
        True if successful, False otherwise.
    """
    config = json.loads(Path(config_path).read_text(encoding="utf-8"))
    slide_style = style or config.get("style", DEFAULT_STYLE)

    # Find slide by id
    slide = None
    for s in config["slides"]:
        if s.get("id") == slide_id:
            slide = s
            break

    if slide is None:
        print(f"Error: slide '{slide_id}' not found in config")
        return False

    img_path = Path(output_dir) / f"{slide_id}.png"

    # Delete existing
    if img_path.exists():
        img_path.unlink()
        print(f"Deleted existing {img_path}")

    # Regenerate
    print(f"Regenerating {slide_id} (style: {slide_style})...")
    ok = generate_slide_image(slide["prompt"], img_path, style=slide_style)

    if not ok:
        print(f"Failed to regenerate {slide_id}")
        return False

    print(f"Regenerated {slide_id} successfully")

    # Rebuild PPTX + HTML preview
    pptx_path = Path(output_dir).parent / "presentation.pptx"
    create_pptx(output_dir, str(pptx_path))
    html_path = Path(output_dir).parent / "preview.html"
    create_html_preview(output_dir, str(html_path))

    return True


def create_pptx_with_notes(image_dir: str, output_path: str, notes_path: str,
                           slide_width_inches: float = 13.333, slide_height_inches: float = 7.5):
    """Package slide images into PPTX with speaker notes.

    Args:
        image_dir: Directory containing slide_*.png files.
        output_path: Output PPTX file path.
        notes_path: Path to JSON file with speaker notes.
                    Format: [{"index": 0, "notes": "Speaker text..."}, ...]
        slide_width_inches: Slide width (default 16:9 widescreen).
        slide_height_inches: Slide height.
    """
    from pptx import Presentation
    from pptx.util import Inches, Emu

    img_dir = Path(image_dir)
    images = sorted(img_dir.glob("slide_*.png"))

    if not images:
        print("No slide images found!")
        return False

    # Load speaker notes
    notes = {}
    notes_data = json.loads(Path(notes_path).read_text(encoding="utf-8"))
    for n in notes_data:
        notes[n["index"]] = n["notes"]

    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)

    blank_layout = prs.slide_layouts[6]

    for i, img_path in enumerate(images):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )

        # Add speaker notes if available for this slide index
        if i in notes:
            slide.notes_slide.notes_text_frame.text = notes[i]

    prs.save(str(output_path))
    print(f"PPTX with notes saved: {output_path} ({len(images)} slides, {len(notes)} notes)")
    return True


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    cmd = sys.argv[1]

    if cmd == "generate":
        if len(sys.argv) < 4:
            print("Usage: whiteboard_generator.py generate <config.json> <output_dir> [style]")
            sys.exit(1)
        config_path = sys.argv[2]
        output_dir = sys.argv[3]
        style = sys.argv[4] if len(sys.argv) > 4 else DEFAULT_STYLE

        paths = generate_from_config(config_path, output_dir, style)

        # Auto-create PPTX and HTML preview
        success = sum(1 for p in paths if p)
        if success > 0:
            pptx_path = Path(output_dir).parent / "presentation.pptx"
            create_pptx(output_dir, str(pptx_path))
            html_path = Path(output_dir).parent / "preview.html"
            create_html_preview(output_dir, str(html_path))

    elif cmd == "test":
        if len(sys.argv) < 4:
            print("Usage: whiteboard_generator.py test <prompt_text> <output_dir> [style]")
            sys.exit(1)
        prompt = sys.argv[2]
        output_dir = Path(sys.argv[3])
        style = sys.argv[4] if len(sys.argv) > 4 else DEFAULT_STYLE
        output_dir.mkdir(parents=True, exist_ok=True)
        img_path = output_dir / "test_slide.png"

        print(f"Generating test slide (style: {style})...")
        ok = generate_slide_image(prompt, img_path, style=style)
        if ok:
            print(f"Success! Image saved to {img_path}")
        else:
            print("Generation failed.")
            sys.exit(1)

    elif cmd == "pptx":
        if len(sys.argv) < 4:
            print("Usage: whiteboard_generator.py pptx <image_dir> <output.pptx>")
            sys.exit(1)
        create_pptx(sys.argv[2], sys.argv[3])

    elif cmd == "html":
        if len(sys.argv) < 4:
            print("Usage: whiteboard_generator.py html <image_dir> <output.html>")
            sys.exit(1)
        create_html_preview(sys.argv[2], sys.argv[3])

    elif cmd == "styles":
        print(f"Available AI styles ({len(STYLES)} total):\n")
        for category, style_names in STYLE_CATEGORIES.items():
            print(f"  {category}:")
            for name in style_names:
                first_line = STYLES[name].strip().split('\n')[0]
                print(f"    {name:20s} {first_line[:60]}...")
            print()
        print(f"  HTML-only styles (use slide_templates.py):")
        for name in HTML_ONLY_STYLES:
            print(f"    {name}")
        print(f"\nAliases:")
        for alias, target in sorted(STYLE_ALIASES.items()):
            print(f"    {alias:20s} -> {target}")

    elif cmd == "regenerate":
        if len(sys.argv) < 5:
            print("Usage: whiteboard_generator.py regenerate <config.json> <output_dir> <slide_id> [style]")
            sys.exit(1)
        config_path = sys.argv[2]
        output_dir = sys.argv[3]
        slide_id = sys.argv[4]
        style = sys.argv[5] if len(sys.argv) > 5 else DEFAULT_STYLE
        ok = regenerate_slide(config_path, output_dir, slide_id, style)
        if not ok:
            sys.exit(1)

    elif cmd == "notes-pptx":
        if len(sys.argv) < 5:
            print("Usage: whiteboard_generator.py notes-pptx <image_dir> <output.pptx> <notes.json>")
            sys.exit(1)
        create_pptx_with_notes(sys.argv[2], sys.argv[3], sys.argv[4])

    elif cmd == "recommend":
        if len(sys.argv) < 3:
            print("Usage: whiteboard_generator.py recommend <task_description>")
            sys.exit(1)
        task = " ".join(sys.argv[2:])
        recs = recommend_styles(task)
        print(f"Recommended styles for '{task}':")
        for i, s in enumerate(recs, 1):
            first_line = STYLES[s].strip().split('\n')[0]
            print(f"  {i}. {s:20s} {first_line[:60]}...")

    else:
        print(f"Unknown command: {cmd}")
        print(__doc__)
        sys.exit(1)


if __name__ == "__main__":
    main()
