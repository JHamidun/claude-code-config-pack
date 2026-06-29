"""
Slide Export - Convert HTML slides to PDF/PPTX via Playwright screenshots.

Adapted from Manus AI slides_to_pdf Go module.
Uses Playwright for headless browser screenshots, Pillow for PDF, python-pptx for PPTX.

Usage:
    python slide_export.py screenshots ./project
    python slide_export.py pdf ./project ./output.pdf
    python slide_export.py pptx ./project ./output.pptx
    python slide_export.py html ./project ./output.html

Dependencies:
    pip install playwright Pillow python-pptx
    playwright install chromium
"""

import json
import sys
import os
import asyncio
from pathlib import Path


async def take_slide_screenshots(project_dir: str, width: int = 1280, height: int = 720) -> list:
    """
    Take screenshots of each active slide using Playwright.

    Returns list of screenshot paths.
    """
    from playwright.async_api import async_playwright

    project_path = Path(project_dir)
    screenshots_dir = project_path / "screenshots"
    screenshots_dir.mkdir(exist_ok=True)

    # Read state to get active slides in order
    state_path = project_path / "slide_state.json"
    if not state_path.exists():
        raise FileNotFoundError(f"slide_state.json not found in {project_dir}")

    with open(state_path, "r", encoding="utf-8") as f:
        state = json.load(f)

    active_slides = []
    for entry in state["outline"]:
        sid = entry["id"]
        if state.get("states", {}).get(sid) == "deleted":
            continue
        slide_path = project_path / "slides" / f"{sid}.html"
        if slide_path.exists():
            active_slides.append((sid, slide_path))

    if not active_slides:
        raise ValueError("No active slides found")

    screenshot_paths = []

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": width, "height": height})

        for i, (sid, slide_path) in enumerate(active_slides):
            file_url = slide_path.resolve().as_uri()
            await page.goto(file_url)
            await page.wait_for_timeout(500)  # Wait for Tailwind CDN + fonts

            screenshot_path = screenshots_dir / f"{i+1:03d}_{sid}.png"
            await page.screenshot(path=str(screenshot_path), full_page=False)
            screenshot_paths.append(str(screenshot_path))

        await browser.close()

    return screenshot_paths


def create_pdf_from_images(image_paths: list, output_path: str, width: int = 1280, height: int = 720) -> str:
    """Create PDF from screenshot images using Pillow."""
    from PIL import Image

    if not image_paths:
        raise ValueError("No images provided")

    images = []
    for path in image_paths:
        img = Image.open(path).convert("RGB")
        # Resize to exact dimensions if needed
        if img.size != (width, height):
            img = img.resize((width, height), Image.LANCZOS)
        images.append(img)

    # Save as PDF (first image, append rest)
    images[0].save(
        output_path,
        "PDF",
        resolution=100.0,
        save_all=True,
        append_images=images[1:] if len(images) > 1 else [],
    )

    return output_path


def create_pptx_from_images(image_paths: list, output_path: str, title: str = "Presentation") -> str:
    """Create PPTX from screenshot images using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Emu

    prs = Presentation()
    # Set slide dimensions to 16:9 (standard widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        # Add image covering full slide
        slide.shapes.add_picture(
            path,
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output_path)
    return output_path


def create_combined_html(project_dir: str, output_path: str) -> str:
    """
    Create a single HTML file with all slides combined.
    Includes keyboard navigation (arrow keys).
    """
    project_path = Path(project_dir)
    state_path = project_path / "slide_state.json"

    with open(state_path, "r", encoding="utf-8") as f:
        state = json.load(f)

    slides_html = []
    for entry in state["outline"]:
        sid = entry["id"]
        if state.get("states", {}).get(sid) == "deleted":
            continue

        slide_path = project_path / "slides" / f"{sid}.html"
        if slide_path.exists():
            with open(slide_path, "r", encoding="utf-8") as f:
                content = f.read()

            # Extract body content between <body> and </body>
            import re
            body_match = re.search(r'<body[^>]*>(.*?)</body>', content, re.DOTALL)
            if body_match:
                slides_html.append(body_match.group(1).strip())
            else:
                slides_html.append(content)

    title = state.get("title", "Presentation")
    total = len(slides_html)

    combined = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<script src="https://cdn.tailwindcss.com"></script>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Playfair+Display:wght@400;700;900&display=swap" rel="stylesheet">
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{ font-family: 'Inter', sans-serif; background: #000; overflow: hidden; }}
  .slide-container {{ display: none; width: 100vw; height: 100vh; }}
  .slide-container.active {{ display: flex; align-items: center; justify-content: center; }}
  .slide {{ width: 1280px; height: 720px; transform-origin: center; position: relative; overflow: hidden; }}
  .progress {{ position: fixed; bottom: 0; left: 0; height: 3px; background: #6366f1; transition: width 0.3s; z-index: 100; }}
  .counter {{ position: fixed; bottom: 12px; right: 20px; color: #666; font-size: 14px; z-index: 100; }}

  @media (max-width: 1280px) or (max-height: 720px) {{
    .slide {{ transform: scale(calc(min(100vw / 1280, 100vh / 720))); }}
  }}
</style>
</head>
<body>
"""

    for i, html in enumerate(slides_html):
        active = ' active' if i == 0 else ''
        combined += f'<div class="slide-container{active}" data-index="{i}">{html}</div>\n'

    combined += f"""
<div class="progress" id="progress"></div>
<div class="counter" id="counter">1 / {total}</div>

<script>
const slides = document.querySelectorAll('.slide-container');
let current = 0;
const total = slides.length;

function goTo(n) {{
  if (n < 0 || n >= total) return;
  slides[current].classList.remove('active');
  current = n;
  slides[current].classList.add('active');
  document.getElementById('counter').textContent = (current + 1) + ' / ' + total;
  document.getElementById('progress').style.width = ((current + 1) / total * 100) + '%';
}}

document.addEventListener('keydown', e => {{
  if (e.key === 'ArrowRight' || e.key === ' ') goTo(current + 1);
  else if (e.key === 'ArrowLeft') goTo(current - 1);
  else if (e.key === 'Home') goTo(0);
  else if (e.key === 'End') goTo(total - 1);
}});

document.addEventListener('click', e => {{
  if (e.clientX > window.innerWidth / 2) goTo(current + 1);
  else goTo(current - 1);
}});

// Init progress
document.getElementById('progress').style.width = (1 / total * 100) + '%';
</script>
</body>
</html>"""

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(combined)

    return output_path


async def export_to_pdf(project_dir: str, output_path: str) -> dict:
    """Full pipeline: screenshots -> PDF."""
    screenshots = await take_slide_screenshots(project_dir)
    create_pdf_from_images(screenshots, output_path)
    return {
        "success": True,
        "output": output_path,
        "slides": len(screenshots),
        "format": "pdf",
    }


async def export_to_pptx(project_dir: str, output_path: str, title: str = "") -> dict:
    """Full pipeline: screenshots -> PPTX."""
    screenshots = await take_slide_screenshots(project_dir)
    if not title:
        state_path = Path(project_dir) / "slide_state.json"
        with open(state_path, "r", encoding="utf-8") as f:
            title = json.load(f).get("title", "Presentation")
    create_pptx_from_images(screenshots, output_path, title)
    return {
        "success": True,
        "output": output_path,
        "slides": len(screenshots),
        "format": "pptx",
    }


def export_to_html(project_dir: str, output_path: str) -> dict:
    """Combine slides into single navigable HTML."""
    create_combined_html(project_dir, output_path)
    state_path = Path(project_dir) / "slide_state.json"
    with open(state_path, "r", encoding="utf-8") as f:
        state = json.load(f)
    active = sum(1 for e in state["outline"] if state.get("states", {}).get(e["id"]) != "deleted")
    return {
        "success": True,
        "output": output_path,
        "slides": active,
        "format": "html",
    }


# --- CLI ---

def main():
    if len(sys.argv) < 3:
        print(json.dumps({
            "error": "Usage: slide_export.py <command> <project_dir> [output_path]",
            "commands": ["screenshots", "pdf", "pptx", "html"],
        }))
        sys.exit(1)

    command = sys.argv[1]
    project_dir = sys.argv[2]

    try:
        if command == "screenshots":
            paths = asyncio.run(take_slide_screenshots(project_dir))
            print(json.dumps({"success": True, "screenshots": paths}))

        elif command == "pdf":
            output = sys.argv[3] if len(sys.argv) > 3 else str(Path(project_dir) / "presentation.pdf")
            result = asyncio.run(export_to_pdf(project_dir, output))
            print(json.dumps(result))

        elif command == "pptx":
            output = sys.argv[3] if len(sys.argv) > 3 else str(Path(project_dir) / "presentation.pptx")
            result = asyncio.run(export_to_pptx(project_dir, output))
            print(json.dumps(result))

        elif command == "html":
            output = sys.argv[3] if len(sys.argv) > 3 else str(Path(project_dir) / "presentation.html")
            result = export_to_html(project_dir, output)
            print(json.dumps(result))

        else:
            print(json.dumps({"error": f"Unknown command: {command}"}))
            sys.exit(1)

    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)


if __name__ == "__main__":
    main()
