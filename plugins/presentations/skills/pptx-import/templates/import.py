#!/usr/bin/env python3
import sys, json, os, base64
from pptx import Presentation
from pptx.util import Emu

EMU_PER_PX = 9525

def emu_px(emu): return int(round(emu / EMU_PER_PX))

def shape_to_dict(shape):
    out = {
        'kind': str(shape.shape_type).split('.')[-1] if shape.shape_type else 'unknown',
        'x': emu_px(shape.left or 0),
        'y': emu_px(shape.top or 0),
        'w': emu_px(shape.width or 0),
        'h': emu_px(shape.height or 0),
    }
    if shape.has_text_frame:
        out['kind'] = 'text'
        out['paragraphs'] = []
        for p in shape.text_frame.paragraphs:
            runs = []
            for r in p.runs:
                runs.append({
                    'text': r.text,
                    'bold': r.font.bold,
                    'italic': r.font.italic,
                    'size': r.font.size.pt if r.font.size else None,
                    'name': r.font.name,
                    'color': str(r.font.color.rgb) if r.font.color and r.font.color.type else None,
                })
            out['paragraphs'].append({ 'runs': runs, 'align': str(p.alignment) if p.alignment else None })
    elif shape.shape_type == 13:  # PICTURE
        out['kind'] = 'image'
        try:
            img = shape.image
            out['ext'] = img.ext
            out['data_b64'] = base64.b64encode(img.blob).decode()
        except Exception: pass
    return out

def main():
    if len(sys.argv) < 2:
        print("Usage: python import.py <file.pptx> [out.json]"); sys.exit(1)
    src = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else src.replace('.pptx', '.json')

    prs = Presentation(src)
    W, H = emu_px(prs.slide_width), emu_px(prs.slide_height)

    slides = []
    for i, sl in enumerate(prs.slides):
        items = [shape_to_dict(sh) for sh in sl.shapes]
        slides.append({ 'index': i, 'items': items })

    with open(out, 'w') as f:
        json.dump({ 'width': W, 'height': H, 'slides': slides }, f, indent=2, default=str)
    print(f"✓ {out} — {len(slides)} slides, {W}×{H}")

if __name__ == "__main__":
    main()
