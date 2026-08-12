#!/usr/bin/env python3
import json, sys, os, urllib.request, tempfile
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

EMU_PER_PX = 9525  # PowerPoint

def px(v): return Emu(int(round(v * EMU_PER_PX)))

def hex_to_rgb(s):
    s = s.lstrip('#')
    return RGBColor(int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))

def fetch_image(src, tmp):
    if src.startswith('data:'):
        import base64
        head, b64 = src.split(',', 1)
        ext = 'png' if 'png' in head else 'jpg'
        path = os.path.join(tmp, f"img_{abs(hash(src))}.{ext}")
        with open(path, 'wb') as f: f.write(base64.b64decode(b64))
        return path
    elif src.startswith('http'):
        path = os.path.join(tmp, f"img_{abs(hash(src))}")
        urllib.request.urlretrieve(src, path)
        return path
    elif src.startswith('file://'):
        return src.replace('file://','')
    else:
        return src  # относительный путь

def main():
    if len(sys.argv) < 2:
        print("Usage: build_pptx.py slides.json [out.pptx]"); sys.exit(1)

    data = json.load(open(sys.argv[1]))
    out = sys.argv[2] if len(sys.argv) > 2 else sys.argv[1].replace('.json', '.pptx')
    W, H = data['width'], data['height']

    prs = Presentation()
    prs.slide_width  = px(W)
    prs.slide_height = px(H)
    blank = prs.slide_layouts[6]

    tmp = tempfile.mkdtemp()

    for s in data['slides']:
        slide = prs.slides.add_slide(blank)
        # Фон
        if s.get('bg'):
            try:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = hex_to_rgb(s['bg'])
            except Exception: pass

        for item in s['items']:
            x,y,w,h = px(item['x']), px(item['y']), px(item['w']), px(item['h'])

            if item['kind'] == 'rect':
                shp = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.RECTANGLE
                shp.fill.solid()
                if item.get('fill'): shp.fill.fore_color.rgb = hex_to_rgb(item['fill'])
                shp.line.fill.background()

            elif item['kind'] == 'text':
                tb = slide.shapes.add_textbox(x, y, w, h)
                tf = tb.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                p = tf.paragraphs[0]
                p.text = item['text']
                run = p.runs[0]
                run.font.name = item.get('fontFamily') or 'Helvetica'
                run.font.size = Pt(item.get('fontSize', 16) * 0.75)  # px → pt approx
                if item.get('color'): run.font.color.rgb = hex_to_rgb(item['color'])
                fw = str(item.get('fontWeight','400'))
                run.font.bold = fw in ('bold','600','700','800','900')
                run.font.italic = bool(item.get('italic'))
                a = (item.get('align') or 'left')
                p.alignment = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT,'justify':PP_ALIGN.JUSTIFY}.get(a, PP_ALIGN.LEFT)

            elif item['kind'] == 'image':
                try:
                    p = fetch_image(item['src'], tmp)
                    slide.shapes.add_picture(p, x, y, w, h)
                except Exception as e:
                    print("img fail:", item['src'], e)

    prs.save(out)
    print("✓", out)

if __name__ == "__main__":
    main()
